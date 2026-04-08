"""
PDF to Office converters — extracts content from PDFs into Word, Excel, and PowerPoint.

Extraction Strategies:
  1. PDF → Word: Uses `pdf2docx` to parse PDF layout elements (paragraphs, tables,
     images) and reconstructs them as a .docx preserving formatting.
  2. PDF → Excel: Uses `pdfplumber` to detect table boundaries via line/string
     alignment heuristics, then extracts structured row/column data into .xlsx.
  3. PDF → PowerPoint: Uses PyMuPDF to extract text blocks and images per page,
     maps their coordinates to a 16:9 slide, and builds a .pptx.

Limitations & Expectations:
  - PDF → Word: Best for text-heavy docs. Scanned PDFs yield empty output.
  - PDF → Excel: Only extracts *tabular* data. Free-form text is ignored.
    If no tables are detected, a clear error is returned.
  - PDF → PowerPoint: Produces an extraction, not a visual clone. Text blocks
    and images are placed roughly where they appear, but complex layouts,
    backgrounds, and fonts will differ from the original PDF.
"""

import logging
from pathlib import Path

from backend.converters import ConversionError

logger = logging.getLogger(__name__)

# ─── Import Guards ──────────────────────────────────────────────────────────
try:
    from pdf2docx import Converter
    _HAS_PDF2DOCX = True
except ImportError:
    _HAS_PDF2DOCX = False
    Converter = None

try:
    import pdfplumber
    _HAS_PDFPLUMBER = True
except ImportError:
    _HAS_PDFPLUMBER = False
    pdfplumber = None

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    _HAS_PPTX = True
except ImportError:
    _HAS_PPTX = False
    Presentation = None

try:
    import fitz
    _HAS_FITZ = True
except ImportError:
    _HAS_FITZ = False
    fitz = None

try:
    import openpyxl
    from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
    _HAS_OPENPYXL = True
except ImportError:
    _HAS_OPENPYXL = False
    openpyxl = None


# ═══════════════════════════════════════════════════════════════════════════════
# 1. PDF to Word
# ═══════════════════════════════════════════════════════════════════════════════

def convert_pdf_to_word(input_path: Path, output_path: Path) -> Path:
    """
    Convert a PDF file to a Word (.docx) file preserving text, tables, and images.

    Args:
        input_path: Path to the source PDF.
        output_path: Path where the .docx will be written.

    Returns:
        Path to the created .docx file.

    Raises:
        ConversionError: If library missing, PDF invalid, encrypted, or empty output.
    """
    if not _HAS_PDF2DOCX:
        raise ConversionError(
            "PDF to Word conversion is not available on this server.",
            detail="pdf2docx is not installed.",
        )

    _validate_pdf(input_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    logger.info("PDF→DOCX: %s (%.1f KB)", input_path.name, input_path.stat().st_size / 1024)

    try:
        with Converter(str(input_path)) as cv:
            cv.convert(str(output_path))
    except ValueError as e:
        error_msg = str(e).lower()
        if "encrypted" in error_msg or "password" in error_msg:
            raise ConversionError(
                "This PDF is encrypted or password-protected.",
                detail=str(e),
            )
        raise ConversionError("Failed to read the PDF.", detail=str(e))
    except RuntimeError as e:
        error_msg = str(e).lower()
        if "empty" in error_msg:
            raise ConversionError(
                "This PDF appears to be scanned (image-only). No text to extract.",
                detail=str(e),
            )
        raise ConversionError("Failed to convert the PDF.", detail=str(e))
    except Exception as e:
        logger.error("pdf2docx unexpected error: %s", e, exc_info=True)
        raise ConversionError("Unexpected error during PDF to Word conversion.", detail=str(e))

    _validate_output(output_path, "Word document")
    logger.info("PDF→DOCX complete: %s (%.1f KB)", output_path.name, output_path.stat().st_size / 1024)
    return output_path


# ═══════════════════════════════════════════════════════════════════════════════
# 2. PDF to Excel
# ═══════════════════════════════════════════════════════════════════════════════

def convert_pdf_to_excel(input_path: Path, output_path: Path) -> tuple[Path, int, int]:
    """
    Extract tabular data from a PDF into an Excel (.xlsx) file.

    Uses pdfplumber to detect table boundaries on each page. Each detected
    table is written to a separate sheet named "Page X - Table Y".

    Args:
        input_path: Path to the source PDF.
        output_path: Path where the .xlsx will be written.

    Returns:
        Tuple of (output_path, total_tables_found, total_rows_extracted).

    Raises:
        ConversionError: If library missing, PDF invalid, or no tables found.
    """
    if not _HAS_PDFPLUMBER or not _HAS_OPENPYXL:
        raise ConversionError(
            "PDF to Excel conversion is not available.",
            detail="pdfplumber or openpyxl is not installed.",
        )

    _validate_pdf(input_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    logger.info("PDF→XLSX: %s (%.1f KB)", input_path.name, input_path.stat().st_size / 1024)

    wb = openpyxl.Workbook()
    # Remove default sheet created by openpyxl
    wb.remove(wb.active)

    total_tables = 0
    total_rows = 0

    try:
        with pdfplumber.open(str(input_path)) as pdf:
            if len(pdf.pages) == 0:
                raise ConversionError("The PDF has 0 pages.")

            for page_idx, page in enumerate(pdf.pages):
                tables = page.extract_tables({
                    "vertical_strategy": "lines",       # Explicit table borders
                    "horizontal_strategy": "lines",     # Explicit table borders
                    "snap_tolerance": 5,                # Snap to nearby lines
                    "join_tolerance": 5,
                })

                # Fallback: if no bordered tables found, try text-based detection
                if not tables or all(not any(cell for row in t for cell in row if cell) for t in tables):
                    tables = page.extract_tables({
                        "vertical_strategy": "text",    # Infer columns from text alignment
                        "horizontal_strategy": "text",  # Infer rows from text gaps
                        "snap_tolerance": 5,
                        "join_tolerance": 5,
                    })

                if not tables:
                    continue

                for table_idx, table_data in enumerate(tables):
                    # Filter out completely empty rows
                    cleaned_table = [
                        [cell.strip() if cell else "" for cell in row]
                        for row in table_data
                        if any(cell and cell.strip() for cell in row)
                    ]

                    if not cleaned_table:
                        continue

                    # Create sheet (Excel limits sheet name to 31 chars)
                    if len(pdf.pages) == 1 and len(tables) == 1:
                        sheet_name = "Extracted Data"
                    else:
                        sheet_name = f"Pg {page_idx + 1} T{table_idx + 1}"
                    
                    ws = wb.create_sheet(title=sheet_name)
                    
                    # Write data
                    for r_idx, row in enumerate(cleaned_table, start=1):
                        for c_idx, cell_val in enumerate(row, start=1):
                            cell = ws.cell(row=r_idx, column=c_idx, value=cell_val)
                            cell.alignment = Alignment(vertical="top", wrap_text=True)

                    # Format header row if detected
                    _format_excel_header(ws, cleaned_table)

                    # Auto-size columns
                    _auto_size_columns(ws)

                    total_tables += 1
                    total_rows += len(cleaned_table)

    except ConversionError:
        raise
    except Exception as e:
        logger.error("PDF→XLSX extraction failed: %s", e, exc_info=True)
        raise ConversionError("Failed to extract tables from the PDF.", detail=str(e))

    if total_tables == 0:
        raise ConversionError(
            "No tables were detected in this PDF. "
            "This tool extracts structured tabular data only. "
            "If the document contains text but no tables, try PDF to Word instead.",
        )

    try:
        wb.save(str(output_path))
    except Exception as e:
        raise ConversionError("Failed to save the Excel file.", detail=str(e))

    logger.info(
        "PDF→XLSX complete: %d tables, %d rows → %s (%.1f KB)",
        total_tables, total_rows, output_path.name, output_path.stat().st_size / 1024
    )
    return output_path, total_tables, total_rows


# ═══════════════════════════════════════════════════════════════════════════════
# 3. PDF to PowerPoint
# ═══════════════════════════════════════════════════════════════════════════════

def convert_pdf_to_powerpoint(input_path: Path, output_path: Path) -> tuple[Path, int]:
    """
    Extract text blocks and images from a PDF into a PowerPoint (.pptx) file.

    Creates one slide per page. Text blocks are positioned based on their
    coordinates in the PDF, scaled to fit a standard 16:9 slide layout.

    Note: This produces an extraction, not a visual clone. Complex layouts
    may not map perfectly to slides.

    Args:
        input_path: Path to the source PDF.
        output_path: Path where the .pptx will be written.

    Returns:
        Tuple of (output_path, slide_count).

    Raises:
        ConversionError: If libraries missing, PDF invalid, or extraction fails.
    """
    if not _HAS_FITZ or not _HAS_PPTX:
        raise ConversionError(
            "PDF to PowerPoint conversion is not available.",
            detail="PyMuPDF or python-pptx is not installed.",
        )

    _validate_pdf(input_path)
    output_path.parent.mkdir(parents=True, exist_ok=True)

    logger.info("PDF→PPTX: %s (%.1f KB)", input_path.name, input_path.stat().st_size / 1024)

    doc = fitz.open(str(input_path))
    slide_count = len(doc)

    if slide_count == 0:
        doc.close()
        raise ConversionError("The PDF has 0 pages.")

    try:
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # Standard 16:9
        prs.slide_height = Inches(7.5)
        
        blank_layout = prs.slide_layouts[6]  # Blank layout

        for page_idx in range(slide_count):
            page = doc.load_page(page_idx)
            slide = prs.slides.add_slide(blank_layout)

            # Calculate scale factors from PDF points to slide inches
            page_w_inches = page.rect.width / 72.0
            page_h_inches = page.rect.height / 72.0
            
            scale_x = (prs.slide_width / 914400.0) / page_w_inches  # EMUs to Inches
            scale_y = (prs.slide_height / 914400.0) / page_h_inches

            # Extract and place text blocks
            blocks = page.get_text("blocks")
            for block in blocks:
                # block format: (x0, y0, x1, y1, text, block_no, block_type)
                if len(block) < 7:
                    continue
                
                x0, y0, x1, y1, text, block_no, block_type = block[0], block[1], block[2], block[3], block[4], block[5], block[6]
                
                if not text or not text.strip():
                    continue

                # Skip image blocks (block_type == 1)
                if block_type == 1:
                    continue

                try:
                    left = int(x0 / 72.0 * scale_x * 914400)
                    top = int(y0 / 72.0 * scale_y * 914400)
                    width = int((x1 - x0) / 72.0 * scale_x * 914400)
                    height = int((y1 - y0) / 72.0 * scale_y * 914400)

                    # Ensure minimum dimensions
                    width = max(width, Emu(914400))  # Min 1 inch
                    height = max(height, Emu(457200)) # Min 0.5 inch

                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    p = tf.paragraphs[0]
                    p.text = text.strip()
                    p.font.size = Pt(11)
                    p.font.name = "Calibri"
                except Exception as e:
                    logger.debug("Skipped text block %s on page %d: %s", block_no, page_idx + 1, e)
                    continue

            # Extract and place images
            image_list = page.get_images(full=True)
            for img_idx, img in enumerate(image_list):
                xref = img[0]
                try:
                    base_image = doc.extract_image(xref)
                    if not base_image or not base_image.get("image"):
                        continue
                    
                    image_bytes = base_image["image"]
                    ext = base_image.get("ext", "png")
                    
                    # Write temp image file (python-pptx requires file paths)
                    temp_img_path = output_path.parent / f"temp_pptx_img_{page_idx}_{img_idx}.{ext}"
                    with open(temp_img_path, "wb") as f:
                        f.write(image_bytes)

                    # Get image bounding boxes
                    rects = page.get_image_rects(xref)
                    if rects:
                        rect = rects[0]
                        left = int(rect.x0 / 72.0 * scale_x * 914400)
                        top = int(rect.y0 / 72.0 * scale_y * 914400)
                        width = int(rect.width / 72.0 * scale_x * 914400)
                        height = int(rect.height / 72.0 * scale_y * 914400)
                        
                        # Prevent zero-dimension crashes
                        if width > 0 and height > 0:
                            slide.shapes.add_picture(str(temp_img_path), left, top, width, height)
                            
                except Exception as e:
                    logger.debug("Skipped image %d on page %d: %s", img_idx, page_idx + 1, e)
                finally:
                    # Always clean up temp image
                    if 'temp_img_path' in dir() and temp_img_path.exists():
                        temp_img_path.unlink(missing_ok=True)

        doc.close()
        prs.save(str(output_path))

    except ConversionError:
        doc.close()
        raise
    except Exception as e:
        doc.close()
        logger.error("PDF→PPTX failed: %s", e, exc_info=True)
        raise ConversionError("Failed to convert PDF to PowerPoint.", detail=str(e))

    _validate_output(output_path, "PowerPoint file")
    logger.info("PDF→PPTX complete: %d slides → %s (%.1f KB)", slide_count, output_path.name, output_path.stat().st_size / 1024)
    return output_path, slide_count


# ═══════════════════════════════════════════════════════════════════════════════
# Internal Helpers
# ═══════════════════════════════════════════════════════════════════════════════

def _validate_pdf(path: Path) -> None:
    """Standard PDF validation before processing."""
    if not path.exists():
        raise ConversionError("The uploaded PDF file was not found.")
    if path.stat().st_size == 0:
        raise ConversionError("The uploaded PDF file is empty (0 bytes).")
    if path.suffix.lower() != ".pdf":
        raise ConversionError("Invalid file type. Expected a PDF file.")


def _validate_output(path: Path, label: str) -> None:
    """Ensure the output file was created and is not empty."""
    if not path.exists() or path.stat().st_size == 0:
        path.unlink(missing_ok=True)
        raise ConversionError(f"The converted {label} is empty. The source PDF may be corrupted.")


def _format_excel_header(ws, table_data: list[list[str]]) -> None:
    """
    Apply header formatting if the first row looks like a header.
    Same heuristic as image_to_excel: non-empty cells, shorter than data rows.
    """
    if len(table_data) < 2 or len(table_data[0]) < 2:
        return

    first_row = table_data[0]
    if not all(cell.strip() for cell in first_row):
        return

    first_avg = sum(len(c) for c in first_row) / len(first_row)
    remaining_lengths = [len(c) for row in table_data[1:] for c in row if c.strip()]
    
    if not remaining_lengths:
        return

    if first_avg >= (sum(remaining_lengths) / len(remaining_lengths)):
        return

    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    header_font = Font(bold=True, color="FFFFFF", size=11)
    
    for col_idx in range(1, len(first_row) + 1):
        cell = ws.cell(row=1, column=col_idx)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)


def _auto_size_columns(ws) -> None:
    """Auto-size column widths based on content length."""
    for col_cells in ws.columns:
        max_length = 0
        col_letter = col_cells[0].column_letter
        for cell in col_cells:
            if cell.value:
                max_length = max(max_length, len(str(cell.value)))
        adjusted_width = min(max(max_length + 3, 8), 50)
        ws.column_dimensions[col_letter].width = adjusted_width
